"""
Generate PowerPoint Dokumentasi - LIGHT THEME (Clean Layout)
Sistem Temu Kembali Informasi: TF-IDF

Run   : python generate_ppt.py
Output: Dokumentasi_IR_TF-IDF.pptx
"""

import io, sys, os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

if hasattr(sys.stdout, 'buffer') and sys.stdout.encoding != 'utf-8':
    sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8', errors='replace')

# ── Palet Warna ──────────────────────────────────────────────
BG     = RGBColor(0xF5, 0xF7, 0xFF)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
NAVY   = RGBColor(0x1A, 0x27, 0x6C)
BLUE   = RGBColor(0x33, 0x5F, 0xE4)
TEAL   = RGBColor(0x00, 0xBF, 0xA5)
ORANGE = RGBColor(0xFF, 0x6D, 0x3B)
PURPLE = RGBColor(0x7C, 0x4D, 0xFF)
BODY   = RGBColor(0x2E, 0x35, 0x60)
MUTED  = RGBColor(0x7A, 0x82, 0xAC)
CARD   = RGBColor(0xFF, 0xFF, 0xFF)
CARD2  = RGBColor(0xEA, 0xEF, 0xFF)
BORDER = RGBColor(0xC8, 0xD2, 0xF0)

# Layout konstanta
W, H = Inches(13.33), Inches(7.5)
HEADER_H   = 1.38   # tinggi zona header
LEFT_W     = 5.20   # lebar kolom kiri (konten)
RIGHT_X    = 5.90   # posisi X kolom kanan (screenshot)
RIGHT_W    = 7.10   # lebar kolom kanan
RIGHT_Y    = 1.48   # posisi Y atas screenshot
RIGHT_H    = 5.72   # tinggi screenshot (besar)
PAD        = 0.65   # padding kiri slide


# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# PRIMITIVES
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━

def set_bg(slide, color=BG):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color

def rect(slide, x, y, w, h, fill, border=None, bpt=0.8):
    s = slide.shapes.add_shape(
        1, Inches(x), Inches(y), Inches(w), Inches(h))
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    if border:
        s.line.color.rgb = border
        s.line.width = Pt(bpt)
    else:
        s.line.fill.background()
    return s

def oval(slide, x, y, w, h, color):
    s = slide.shapes.add_shape(
        9, Inches(x), Inches(y), Inches(w), Inches(h))
    s.fill.solid()
    s.fill.fore_color.rgb = color
    s.line.fill.background()
    return s

def txt(slide, text, x, y, w, h,
        size=12, bold=False, color=BODY,
        align=PP_ALIGN.LEFT, italic=False):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tb.word_wrap = True
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = color
    r.font.italic = italic
    return tb

def card(slide, x, y, w, h, fill=CARD, border=BORDER):
    rect(slide, x + 0.04, y + 0.04, w, h, BORDER)   # shadow
    rect(slide, x, y, w, h, fill, border, 0.7)

def chip(slide, x, y, w, label, color=BLUE):
    rect(slide, x, y, w, 0.30, color)
    txt(slide, label, x + 0.05, y + 0.03, w - 0.1, 0.25,
        size=8, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

def dot(slide, x, y, r=0.34, color=BLUE):
    oval(slide, x, y, r, r, color)

def bullet_list(slide, items, x, y, w, size=11, gap=0.50, color=BODY):
    for i, item in enumerate(items):
        txt(slide, f"\u2022  {item}", x, y + i * gap, w, gap + 0.05,
            size=size, color=color)

def numbered_steps(slide, items, x, y, w, color_list=None, size=11, gap=0.62):
    """items = list of (title, desc) tuples, warna bergantian."""
    colors = color_list or [BLUE, TEAL, ORANGE, PURPLE]
    for i, item in enumerate(items):
        c = colors[i % len(colors)]
        title = item[0] if isinstance(item, (list, tuple)) else item
        desc  = item[1] if isinstance(item, (list, tuple)) and len(item) > 1 else ""
        yy = y + i * gap
        dot(slide, x, yy + 0.02, 0.30, c)
        txt(slide, str(i + 1), x + 0.002, yy + 0.04, 0.30, 0.28,
            size=9, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        txt(slide, title, x + 0.42, yy + 0.02, w - 0.45, 0.28,
            size=size, bold=True, color=NAVY)
        if desc:
            txt(slide, desc, x + 0.42, yy + 0.30, w - 0.45, 0.28,
                size=size - 0.5, color=BODY)

# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# LAYOUT HELPERS
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━

def slide_header(slide, title, subtitle, accent=BLUE):
    """Garis atas + judul + subtitle standar."""
    rect(slide, 0, 0, 13.33, 0.08, accent)          # garis aksen atas
    rect(slide, 0, 0.08, 0.08, 7.42, CARD2)         # sidebar kiri tipis
    txt(slide, title, PAD, 0.18, 12.3, 0.75,
        size=26, bold=True, color=NAVY)
    rect(slide, PAD, 0.95, 12.3, 0.04, accent)      # garis divider
    if subtitle:
        txt(slide, subtitle, PAD, 1.03, 12.3, 0.38,
            size=10.5, color=MUTED, italic=True)

def screenshot_right(slide, label="Screenshot",
                     x=RIGHT_X, y=RIGHT_Y, w=RIGHT_W, h=RIGHT_H):
    """Kotak screenshot besar di sisi kanan."""
    card(slide, x, y, w, h, CARD2, BORDER)
    mid_y = y + h / 2
    txt(slide, "\U0001f4f7", x, mid_y - 0.65, w, 0.55,
        size=28, color=MUTED, align=PP_ALIGN.CENTER)
    txt(slide, "[ Screenshot ]", x, mid_y - 0.08, w, 0.40,
        size=13, bold=True, color=MUTED, align=PP_ALIGN.CENTER)
    txt(slide, label, x, mid_y + 0.36, w, 0.50,
        size=10, color=MUTED, italic=True, align=PP_ALIGN.CENTER)

def left_section_title(slide, label, y, color=NAVY):
    txt(slide, label, PAD, y, LEFT_W, 0.38,
        size=12, bold=True, color=color)
    rect(slide, PAD, y + 0.40, LEFT_W, 0.03, BORDER)

# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# SLIDES
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━

# ── 01  Cover ────────────────────────────────────────────────
def slide_cover(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s, WHITE)

    # Panel kiri solid biru
    rect(s, 0, 0, 5.6, 7.5, BLUE)
    oval(s, 2.5, -0.6, 3.2, 3.2, RGBColor(0x55, 0x80, 0xFF))   # dekorasi
    oval(s, -0.8, 5.2, 3.0, 3.0, RGBColor(0x1A, 0x48, 0xD0))   # dekorasi

    txt(s, "DOKUMENTASI TEKNIS", 0.35, 1.55, 4.9, 0.38,
        size=9, bold=True, color=RGBColor(0xB0, 0xC8, 0xFF))
    txt(s, "Sistem Temu\nKembali Informasi", 0.35, 1.95, 4.9, 1.5,
        size=29, bold=True, color=WHITE)
    txt(s, "Berbasis TF-IDF &\nCosine Similarity", 0.35, 3.58, 4.9, 1.0,
        size=16, color=RGBColor(0xC5, 0xD8, 0xFF))
    rect(s, 0.35, 4.75, 4.0, 0.05, RGBColor(0xB0, 0xC8, 0xFF))
    txt(s, "Mata Kuliah: Sistem Temu Kembali Informasi",
        0.35, 4.88, 4.9, 0.38, size=9.5, color=RGBColor(0xA0, 0xB8, 0xFF))

    # Kanan: daftar konten
    txt(s, "Konten Presentasi", 5.95, 1.3, 7.0, 0.42,
        size=14, bold=True, color=NAVY)
    rect(s, 5.95, 1.76, 7.0, 0.05, BLUE)

    items = [
        (PURPLE, "Struktur File & Folder"),
        (TEAL,   "Alur Program (Pipeline)"),
        (ORANGE, "Dataset Dokumen PDF"),
        (BLUE,   "Preprocessing Teks (NLP)"),
        (TEAL,   "Feature Extraction (TF-IDF)"),
        (ORANGE, "Cosine Similarity & Ranking"),
        (PURPLE, "Auto-Correct Query"),
        (BLUE,   "Antarmuka Web (Flask)"),
    ]
    for i, (color, label) in enumerate(items):
        y = 1.95 + i * 0.60
        oval(s, 5.95, y + 0.05, 0.28, 0.28, color)
        txt(s, str(i + 1), 5.953, y + 0.07, 0.28, 0.25,
            size=8, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        txt(s, label, 6.35, y + 0.04, 6.3, 0.34, size=12, color=BODY)


# ── 02  Struktur Folder ───────────────────────────────────────
def slide_struktur(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s)
    slide_header(s, "Struktur File & Folder",
                 "Organisasi kode program secara keseluruhan", PURPLE)

    # Pohon folder
    card(s, PAD, 1.48, 5.20, 5.75, CARD, BORDER)
    rect(s, PAD, 1.48, 5.20, 0.06, PURPLE)

    tree = [
        (0, NAVY,   True,  "tugas2/"),
        (1, BLUE,   True,  "\u251c\u2500 app.py"),
        (1, BLUE,   True,  "\u251c\u2500 main.py"),
        (1, MUTED,  False, "\u251c\u2500 generate_ppt.py"),
        (1, ORANGE, True,  "\u251c\u2500 docs/"),
        (2, MUTED,  False, "\u2502   \u251c\u2500 Banten.pdf"),
        (2, MUTED,  False, "\u2502   \u251c\u2500 Daerah_Istimewa_Yogyakarta.pdf"),
        (2, MUTED,  False, "\u2502   \u2514\u2500 ... (10 file PDF total)"),
        (1, TEAL,   True,  "\u251c\u2500 modules/"),
        (2, BODY,   False, "\u2502   \u251c\u2500 loader.py"),
        (2, BODY,   False, "\u2502   \u251c\u2500 preprocessing.py"),
        (2, BODY,   False, "\u2502   \u251c\u2500 feature_extraction.py"),
        (2, BODY,   False, "\u2502   \u251c\u2500 similarity.py"),
        (2, BODY,   False, "\u2502   \u2514\u2500 tfidf.py"),
        (1, PURPLE, True,  "\u2514\u2500 templates/"),
        (2, BODY,   False, "    \u2514\u2500 index.html"),
    ]
    for i, (lvl, color, bold_, name) in enumerate(tree):
        y = 1.62 + i * 0.33
        txt(s, name, PAD + 0.18 + lvl * 0.22, y,
            LEFT_W - 0.4, 0.30, size=10.5, bold=bold_, color=color)

    # Kanan: penjelasan komponen
    comps = [
        (BLUE,   "app.py",
         "Entry point Flask. Routing GET/POST, serve PDF."),
        (BLUE,   "main.py",
         "IRSystem: orkestrator semua modul."),
        (ORANGE, "docs/",
         "10 file PDF provinsi, dibaca otomatis saat startup."),
        (TEAL,   "modules/",
         "Modul NLP: loader, preprocessor, TF-IDF, similarity."),
        (PURPLE, "templates/",
         "index.html: UI Jinja2 untuk halaman web pencarian."),
    ]

    for i, (color, name, desc) in enumerate(comps):
        y = 1.48 + i * 1.12
        card(s, RIGHT_X, y, RIGHT_W, 1.03, CARD, BORDER)
        rect(s, RIGHT_X, y, RIGHT_W, 0.06, color)
        oval(s, RIGHT_X + 0.18, y + 0.24, 0.38, 0.38, color)
        txt(s, name, RIGHT_X + 0.70, y + 0.12, 5.8, 0.32,
            size=13, bold=True, color=NAVY)
        txt(s, desc, RIGHT_X + 0.70, y + 0.52, 6.1, 0.38,
            size=11, color=BODY)


# ── 03  Alur Program ─────────────────────────────────────────
def slide_alur(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s)
    slide_header(s, "Alur Program (Pipeline)",
                 "Dari dokumen PDF hingga hasil pencarian terurut", TEAL)

    steps = [
        (BLUE,   "Load PDF",
         "Scan folder docs/, buka tiap PDF,\nextrak teks halaman per halaman."),
        (TEAL,   "Preprocessing",
         "Case fold \u2192 hapus simbol \u2192 tokenize\n\u2192 stopword removal \u2192 Sastrawi stem."),
        (ORANGE, "TF-IDF",
         "Bangun vocabulary, hitung TF, IDF,\nTF-IDF per term per dokumen."),
        (PURPLE, "Query Input",
         "User ketik query \u2192 preprocess\n\u2192 hitung vektor TF-IDF query."),
        (BLUE,   "Cosine Similarity",
         "Hitung similarity query vs.\nsetiap vektor dokumen."),
        (TEAL,   "Ranking & Tampil",
         "Urutkan skor DESC, buat snippet\ndinamis + highlight kata kunci."),
    ]

    # 3 kolom x 2 baris
    for i, (color, label, desc) in enumerate(steps):
        col = i % 3
        row = i // 3
        x = PAD + col * 4.20
        y = 1.52 + row * 2.80

        card(s, x, y, 3.95, 2.58, CARD, BORDER)
        rect(s, x, y, 3.95, 0.07, color)

        # Nomor bulat
        oval(s, x + 0.22, y + 0.22, 0.48, 0.48, color)
        txt(s, str(i + 1), x + 0.22, y + 0.26, 0.48, 0.35,
            size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

        txt(s, label, x + 0.85, y + 0.26, 2.90, 0.36,
            size=13, bold=True, color=NAVY)
        txt(s, desc, x + 0.22, y + 0.82, 3.52, 1.0,
            size=11, color=BODY)

        # panah kanan (dalam row)
        if col < 2:
            txt(s, "\u2192", x + 3.97, y + 1.15, 0.26, 0.36,
                size=18, bold=True, color=MUTED, align=PP_ALIGN.CENTER)

    txt(s,
        "* Tahap 1\u20133 (Load \u2192 Preprocess \u2192 TF-IDF) berjalan SEKALI saat startup. "
        "Tahap 4\u20136 berjalan setiap kali query baru masuk. *",
        PAD, 7.05, 12.4, 0.35, size=9, color=MUTED, italic=True)


# ── 04  Dataset ──────────────────────────────────────────────
def slide_dataset(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s)
    slide_header(s, "Dataset & Cara Baca Dokumen",
                 "Korpus 10 PDF Provinsi Indonesia  |  modules/loader.py", ORANGE)

    # Daftar dokumen kiri
    left_section_title(s, "Daftar Dokumen (folder docs/)", 1.52, ORANGE)
    docs = [
        ("Banten",                     BLUE),
        ("Daerah Istimewa Yogyakarta",  TEAL),
        ("Daerah Khusus Ibukota Jakarta", ORANGE),
        ("Jawa Barat",                 PURPLE),
        ("Jawa Tengah",                BLUE),
        ("Jawa Timur",                 TEAL),
        ("Kalimantan Barat",           ORANGE),
        ("Kalimantan Selatan",         PURPLE),
        ("Kalimantan Tengah",          BLUE),
        ("Kalimantan Timur",           TEAL),
    ]
    for i, (name, color) in enumerate(docs):
        y = 2.05 + i * 0.52
        card(s, PAD, y, LEFT_W, 0.44, CARD, BORDER)
        rect(s, PAD, y, 0.12, 0.44, color)
        txt(s, f"  {i+1:02d}.  {name}", PAD + 0.22, y + 0.08,
            LEFT_W - 0.3, 0.30, size=11, color=BODY)

    # Cara baca kanan – besar
    card(s, RIGHT_X, RIGHT_Y, RIGHT_W, RIGHT_H, CARD, BORDER)
    rect(s, RIGHT_X, RIGHT_Y, RIGHT_W, 0.06, ORANGE)

    txt(s, "Cara Baca PDF  (DocumentLoader)", RIGHT_X + 0.25, RIGHT_Y + 0.12,
        RIGHT_W - 0.4, 0.36, size=13, bold=True, color=NAVY)
    rect(s, RIGHT_X + 0.25, RIGHT_Y + 0.52, RIGHT_W - 0.5, 0.03, BORDER)

    loader_steps = [
        (BLUE,   "Scan Folder",
         "os.listdir() mencari semua file *.pdf\ndalam folder dataset yang dikonfigurasi."),
        (TEAL,   "Buka File",
         "Menggunakan PyMuPDF atau pdfminer\nuntuk membuka file satu per satu."),
        (ORANGE, "Ekstrak Teks",
         "Iterasi halaman demi halaman,\ngabungkan seluruh teks jadi satu string."),
        (PURPLE, "Simpan ke Dict",
         "{ 'nama.pdf': 'isi teks...' }\nDiserahkan ke TextPreprocessor."),
    ]
    for i, (color, ttl, desc) in enumerate(loader_steps):
        y = RIGHT_Y + 0.65 + i * 1.22
        card(s, RIGHT_X + 0.25, y, RIGHT_W - 0.5, 1.08, CARD2, BORDER)
        rect(s, RIGHT_X + 0.25, y, 0.10, 1.08, color)
        txt(s, f"{i+1}.  {ttl}", RIGHT_X + 0.48, y + 0.10,
            RIGHT_W - 0.80, 0.30, size=12, bold=True, color=NAVY)
        txt(s, desc, RIGHT_X + 0.48, y + 0.48,
            RIGHT_W - 0.80, 0.52, size=11, color=BODY)


# ── 05  Preprocessing ────────────────────────────────────────
def slide_preprocessing(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s)
    slide_header(s, "Fitur: Preprocessing Teks",
                 "modules/preprocessing.py  |  class TextPreprocessor", BLUE)

    steps = [
        (BLUE,   "Case Folding",
         "Ubah semua huruf ke huruf kecil.\n\"Jawa BARAT\" \u2192 \"jawa barat\""),
        (TEAL,   "Remove Non-Alpha",
         "Hapus angka & simbol, ganti spasi.\nRegex: [^a-zA-Z\\s]"),
        (ORANGE, "Tokenisasi",
         "Pecah kalimat menjadi daftar kata.\nNLTK word_tokenize()"),
        (PURPLE, "Stopword Removal",
         "Hapus kata tak bermakna (NLTK BI).\n\"dan\", \"yang\", \"di\", \"ke\"..."),
        (BLUE,   "Stemming",
         "Kembalikan kata ke bentuk dasar.\nSastrawi Stemmer (Bahasa Indonesia)"),
    ]

    for i, (color, title, desc) in enumerate(steps):
        y = 1.52 + i * 1.08
        card(s, PAD, y, LEFT_W, 0.96, CARD, BORDER)
        rect(s, PAD, y, 0.10, 0.96, color)
        oval(s, PAD + 0.18, y + 0.28, 0.36, 0.36, color)
        txt(s, str(i+1), PAD + 0.18, y + 0.31, 0.36, 0.28,
            size=9.5, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        txt(s, title, PAD + 0.68, y + 0.10, LEFT_W - 0.78, 0.30,
            size=12, bold=True, color=NAVY)
        txt(s, desc, PAD + 0.68, y + 0.46, LEFT_W - 0.78, 0.46,
            size=11, color=BODY)

    # Screenshot besar kanan
    screenshot_right(s,
        label="Tabel preprocessing dokumen\n(tampilan di web UI)")


# ── 06  TF-IDF ───────────────────────────────────────────────
def slide_tfidf(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s)
    slide_header(s, "Fitur: Perhitungan TF-IDF",
                 "modules/feature_extraction.py  |  class FeatureExtractor", TEAL)

    # Rumus
    left_section_title(s, "Rumus", 1.52, TEAL)
    formulas = [
        (BLUE,   "TF(t, d)",   "= Jumlah kemunculan term t dalam dokumen d"),
        (TEAL,   "IDF(t)",     "= log\u2081\u2080( N \u00f7 df(t) )   \u2014   N = total dokumen"),
        (ORANGE, "TF-IDF",     "= TF(t, d)  \u00d7  IDF(t)   (bobot akhir term)"),
    ]
    for i, (color, label, formula) in enumerate(formulas):
        y = 2.02 + i * 0.56
        card(s, PAD, y, LEFT_W, 0.48, CARD, BORDER)
        chip(s, PAD + 0.12, y + 0.09, 1.05, label, color)
        txt(s, formula, PAD + 1.28, y + 0.10, LEFT_W - 1.42, 0.32,
            size=11, color=BODY)

    # Tahapan kode
    left_section_title(s, "Tahapan di Kode", 3.83, TEAL)
    code_steps = [
        (BLUE,   "compute_tf()",
         "Counter(tokens) per dokumen"),
        (TEAL,   "build_vocabulary()",
         "Kumpulkan unique term, sort alfabet"),
        (ORANGE, "compute_df()",
         "Hitung jumlah dokumen per term"),
        (PURPLE, "compute_idf()",
         "log10(N / df[term]) tiap term"),
        (BLUE,   "compute_tfidf()",
         "tf[doc][term] \u00d7 idf[term]"),
        (TEAL,   "get_vector()",
         "Dict TF-IDF \u2192 list vektor"),
    ]
    for i, (color, func, desc) in enumerate(code_steps):
        y = 4.32 + i * 0.52
        card(s, PAD, y, LEFT_W, 0.44, CARD, BORDER)
        rect(s, PAD, y, 0.08, 0.44, color)
        txt(s, func, PAD + 0.18, y + 0.04, 2.1, 0.26,
            size=10.5, bold=True, color=color)
        txt(s, desc, PAD + 2.35, y + 0.06, LEFT_W - 2.45, 0.30,
            size=10, color=BODY)

    # Screenshot besar kanan
    screenshot_right(s,
        label="Tabel TF-IDF per term query & dokumen\n(tampilan di web UI)")


# ── 07  Similarity ───────────────────────────────────────────
def slide_similarity(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s)
    slide_header(s, "Fitur: Similarity & Ranking",
                 "modules/similarity.py  |  class SimilarityCalculator", ORANGE)

    # Cosine
    left_section_title(s, "Cosine Similarity  (Metode Utama)", 1.52, ORANGE)
    card(s, PAD, 2.0, LEFT_W, 0.80, CARD, BORDER)
    rect(s, PAD, 2.0, LEFT_W, 0.07, ORANGE)
    txt(s, "cos(q, d)  =  (q \u00b7 d)  /  (|q|  \u00d7  |d|)",
        PAD, 2.14, LEFT_W, 0.44, size=14, bold=True,
        color=NAVY, align=PP_ALIGN.CENTER)
    txt(s, "Rentang nilai: 0.0 (tidak relevan) \u2014 1.0 (sangat relevan)",
        PAD, 2.60, LEFT_W, 0.28, size=10, color=MUTED, align=PP_ALIGN.CENTER)

    # Euclidean
    left_section_title(s, "Euclidean Distance  (Alternatif)", 3.0, PURPLE)
    card(s, PAD, 3.48, LEFT_W, 0.72, CARD, BORDER)
    rect(s, PAD, 3.48, LEFT_W, 0.07, PURPLE)
    txt(s, "d(q, d)  =  \u221a \u03a3 ( q\u1d62 \u2212 d\u1d62 )\u00b2",
        PAD, 3.62, LEFT_W, 0.44, size=14, bold=True,
        color=NAVY, align=PP_ALIGN.CENTER)
    txt(s, "Semakin kecil nilai  \u2192  semakin mirip",
        PAD, 4.08, LEFT_W, 0.26, size=10, color=MUTED, align=PP_ALIGN.CENTER)

    # Ranking
    left_section_title(s, "Proses Ranking", 4.42, BLUE)
    rank_steps = [
        "Hitung similarity query vs. semua dokumen (loop)",
        "Simpan list: [{ filename, score, snippet }]",
        "Urutkan DESC berdasarkan score",
        "Assign rank: 1, 2, 3, ... setelah sorting",
    ]
    bullet_list(s, rank_steps, PAD, 4.92, LEFT_W, size=11, gap=0.50)

    # Screenshot besar kanan
    screenshot_right(s,
        label="Tabel ranking hasil pencarian\n& similarity matrix antar dokumen")


# ── 08  Auto-Correct ─────────────────────────────────────────
def slide_autocorrect(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s)
    slide_header(s, "Fitur: Auto-Correct Query",
                 "main.py  |  IRSystem.search()  \u2014  menggunakan difflib", PURPLE)

    left_section_title(s, "Cara Kerja", 1.52, PURPLE)
    steps = [
        ("Bangun valid_words",   "Set semua token dari dokumen"),
        ("Tokenisasi Query",     "Pecah input jadi token-token"),
        ("Cek Tiap Token",       "Ada di valid_words? Ya \u2192 lanjut"),
        ("difflib Match",        "get_close_matches(cutoff=0.75, n=1)"),
        ("Susun Ulang Query",    "Gabung token yang telah terkoreksi"),
        ("Notifikasi UI",        "'Menampilkan hasil untuk: ...'"),
    ]
    numbered_steps(s, steps, PAD, 2.02, LEFT_W,
                   [BLUE, TEAL, ORANGE, PURPLE], size=11, gap=0.70)

    # Contoh
    left_section_title(s, "Contoh Koreksi", 6.52, PURPLE)
    examples = [
        ('"jakrta"',      '"jakarta"',    BLUE),
        ('"klimantan"',   '"kalimantan"', TEAL),
        ('"jogyakarta"',  '"yogyakarta"', ORANGE),
    ]
    ex_y = 7.0
    for inp, out, color in examples:
        txt(s, f"{inp}  \u2192  {out}", PAD, ex_y, LEFT_W, 0.30,
            size=10.5, color=color)
        ex_y += 0.30  # tidak cukup di sini; pindah ke foto

    # Screenshot besar kanan
    screenshot_right(s,
        label="Notifikasi auto-correct\ndi halaman hasil pencarian web")


# ── 09  Snippet & Highlight ───────────────────────────────────
def slide_snippet(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s)
    slide_header(s, "Fitur: Snippet Dinamis & Highlight",
                 "main.py  |  _generate_snippet()  &  _highlight_snippet()", TEAL)

    # Snippet
    left_section_title(s, "Snippet Dinamis", 1.52, TEAL)
    snippet_pts = [
        ("Cari Term",    "Temukan posisi term terpanjang query di teks"),
        ("Potong Teks",  "Ambil window ±150 karakter di sekitar term"),
        ("Ellipsis",     "Tambah '...' jika potongan di tengah teks"),
        ("Fallback",     "Jika term tak ditemukan, ambil 300 char pertama"),
    ]
    numbered_steps(s, snippet_pts, PAD, 2.02, LEFT_W,
                   [TEAL, BLUE, ORANGE, PURPLE], size=11, gap=0.72)

    # Highlight
    left_section_title(s, "Highlight Kata Kunci", 5.05, BLUE)
    hl_pts = [
        ("Escape HTML",   "markupsafe.escape(snippet) agar aman dari XSS"),
        ("Regex Replace", "re.sub case-insensitive untuk tiap token query"),
        ("Wrap <mark>",   "Bungkus term cocok dengan <mark>...</mark>"),
    ]
    numbered_steps(s, hl_pts, PAD, 5.52, LEFT_W,
                   [BLUE, TEAL, ORANGE], size=11, gap=0.66)

    # Screenshot besar kanan
    screenshot_right(s,
        label="Snippet dengan kata kunci ter-highlight\n(tampilan di halaman hasil)")


# ── 10  Web UI ───────────────────────────────────────────────
def slide_webui(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s)
    slide_header(s, "Fitur: Antarmuka Web",
                 "app.py  |  templates/index.html  \u2014  Flask Web Application", BLUE)

    # Fitur
    left_section_title(s, "Fitur Halaman Web", 1.52, BLUE)
    features = [
        "Search bar Google-style, submit via POST",
        "Daftar hasil: rank, skor, nama, snippet + highlight",
        "Notifikasi banner jika query dikoreksi otomatis",
        "Klik nama dokumen \u2192 buka PDF langsung di browser",
        "Tabel preprocessing query (token \u2192 clean \u2192 stem)",
        "Tabel TF-IDF per term query & dokumen",
        "Similarity matrix antar semua dokumen (N\u00d7N)",
    ]
    bullet_list(s, features, PAD, 2.02, LEFT_W, size=11, gap=0.52)

    # Routing
    left_section_title(s, "Routing Flask", 5.82, BLUE)
    routes = [
        (BLUE,   "GET  /",            "Tampilkan halaman kosong"),
        (TEAL,   "POST /",            "Jalankan pencarian, kirim ke template"),
        (ORANGE, "GET /docs/<file>",  "Serve file PDF ke browser"),
    ]
    for i, (color, route, desc) in enumerate(routes):
        y = 6.28 + i * 0.43
        chip(s, PAD, y, 1.85, route, color)
        txt(s, desc, PAD + 1.95, y + 0.02, LEFT_W - 2.0, 0.28,
            size=10.5, color=BODY)

    # 2 screenshot kanan – stacked
    MID = RIGHT_Y + RIGHT_H / 2 + 0.10
    H1  = MID - RIGHT_Y - 0.12
    H2  = RIGHT_Y + RIGHT_H - MID

    screenshot_right(s,
        label="Halaman utama + daftar hasil pencarian",
        x=RIGHT_X, y=RIGHT_Y, w=RIGHT_W, h=H1)

    screenshot_right(s,
        label="Tabel TF-IDF detail & similarity matrix",
        x=RIGHT_X, y=MID, w=RIGHT_W, h=H2)


# ── 11  Penutup ───────────────────────────────────────────────
def slide_penutup(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s, WHITE)

    rect(s, 0, 0, 5.6, 7.5, BLUE)
    oval(s, 2.3, -0.7, 3.4, 3.4, RGBColor(0x55, 0x80, 0xFF))
    oval(s, -1.0, 5.0, 3.2, 3.2, RGBColor(0x1A, 0x48, 0xD0))

    txt(s, "Terima Kasih", 0.3, 2.55, 5.0, 0.95,
        size=38, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    txt(s, "Ada pertanyaan?", 0.3, 3.62, 5.0, 0.48,
        size=16, color=RGBColor(0xC5, 0xD8, 0xFF),
        align=PP_ALIGN.CENTER, italic=True)
    rect(s, 0.6, 4.25, 4.4, 0.05, RGBColor(0xAA, 0xC8, 0xFF))
    txt(s, "Jalankan:  python app.py\nURL:  http://localhost:5000",
        0.3, 4.4, 5.0, 0.70, size=11,
        color=RGBColor(0xA0, 0xB8, 0xFF), align=PP_ALIGN.CENTER)

    # Ringkasan kanan
    txt(s, "Ringkasan Sistem", 5.95, 1.3, 7.0, 0.45,
        size=15, bold=True, color=NAVY)
    rect(s, 5.95, 1.80, 7.0, 0.05, BLUE)

    summary = [
        (BLUE,   "10 Dokumen PDF",  "Profil 10 provinsi Indonesia"),
        (TEAL,   "5 Tahap NLP",     "Case fold \u2192 Tokenize \u2192 Stopword \u2192 Stem"),
        (ORANGE, "TF-IDF Vectors",  "Representasi bobot term per dokumen"),
        (PURPLE, "Cosine Similarity","Metrik relevansi query vs. dokumen"),
        (BLUE,   "Auto-Correct",    "Toleransi typo via difflib (0.75)"),
        (TEAL,   "Flask Web UI",    "Real-time search + tabel detail"),
    ]
    for i, (color, title, desc) in enumerate(summary):
        y = 2.0 + i * 0.87
        card(s, 5.95, y, 7.0, 0.77, CARD, BORDER)
        rect(s, 5.95, y, 0.12, 0.77, color)
        txt(s, title, 6.18, y + 0.08, 2.2, 0.30,
            size=12, bold=True, color=color)
        txt(s, desc, 8.42, y + 0.10, 4.3, 0.52,
            size=11, color=BODY)


# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# MAIN
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━

def main():
    prs = Presentation()
    prs.slide_width  = W
    prs.slide_height = H

    slides = [
        (slide_cover,        "01 - Cover"),
        (slide_struktur,     "02 - Struktur File & Folder"),
        (slide_alur,         "03 - Alur Program"),
        (slide_dataset,      "04 - Dataset"),
        (slide_preprocessing,"05 - Preprocessing"),
        (slide_tfidf,        "06 - TF-IDF"),
        (slide_similarity,   "07 - Similarity & Ranking"),
        (slide_autocorrect,  "08 - Auto-Correct"),
        (slide_snippet,      "09 - Snippet & Highlight"),
        (slide_webui,        "10 - Antarmuka Web"),
        (slide_penutup,      "11 - Penutup"),
    ]

    print("[*] Membuat PPT ...")
    for fn, label in slides:
        print(f"  [OK] Slide {label}")
        fn(prs)

    out = os.path.join(os.path.dirname(os.path.abspath(__file__)),
                       "Dokumentasi_IR_TF-IDF.pptx")
    prs.save(out)
    print(f"\n[DONE] Tersimpan: {out}")
    print("[INFO] Ganti kotak '[Screenshot]' dengan screenshot aktual di PowerPoint.")


if __name__ == "__main__":
    main()
